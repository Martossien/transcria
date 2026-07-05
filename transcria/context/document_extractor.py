"""Extraction déterministe du texte des documents joints (support de réunion).

But : transformer un fichier présenté en réunion (PDF, Word .docx, PowerPoint .pptx,
texte brut) en **texte exploitable** par la LLM de résumé, dans le même canal que
l'invitation (cf. [invite_parser]). Les **images sont ignorées** en v1 (comptées mais
non analysées) — l'analyse visuelle par la LLM viendra plus tard.

Principes (repris de ``invite_parser`` pour cohérence) :

- **Formats XML modernes uniquement**, 100 % pur-Python : ``.pdf`` (pypdf), ``.docx``
  (python-docx), ``.pptx`` (python-pptx), ``.txt``. Les binaires hérités ``.doc``/``.ppt``
  (Office ≤ 2003) ne sont pas gérés — la route renvoie un message « convertissez en
  .docx/.pptx/PDF ».
- **Minimisation PII** : les adresses e-mail sont retirées du texte extrait (même motif
  que l'invitation) ; on ne conserve **jamais** le binaire d'origine, seulement le texte
  assaini et plafonné.
- **Robuste** : un fichier corrompu ou illisible lève ``DocumentExtractionError`` (jamais
  un crash non typé côté route). Un document sans texte (PDF scanné image) rend un texte
  vide et ``images_skipped`` renseigné — pas une erreur.
- **Déterministe** : aucun réseau, aucun GPU, aucune donnée métier en dur.
"""
from __future__ import annotations

import io
import zipfile
from dataclasses import dataclass, field

from transcria.context.invite_parser import _EMAIL_RE, _collapse_ws

# Plafond par défaut du texte conservé par document. Un cran au-dessus de
# ``MAX_BRIEF_CHARS`` (invitation) : un support de réunion porte plus de contexte qu'une
# simple invitation, mais on borne pour ne pas noyer la LLM ni exploser le budget de
# contexte. Surchargé par ``security.max_document_chars`` via la route.
DEFAULT_MAX_CHARS = 12000

# Extensions gérées (XML modernes + texte). En minuscules, avec le point.
SUPPORTED_EXTENSIONS = (".pdf", ".docx", ".pptx", ".txt")


class DocumentExtractionError(Exception):
    """Le document n'a pas pu être lu (format non géré, fichier corrompu, vide)."""


@dataclass
class ExtractedDocument:
    """Résultat d'extraction : texte assaini + métadonnées d'affichage."""

    text: str
    format: str  # "pdf" | "docx" | "pptx" | "txt"
    pages: int = 0  # PDF : pages ; DOCX/TXT : 0 (non pertinent)
    slides: int = 0  # PPTX : diapositives ; sinon 0
    images_skipped: int = 0  # images rencontrées et NON analysées (vision à venir)
    truncated: bool = False  # texte tronqué au plafond
    warnings: list[str] = field(default_factory=list)


def _sanitize(text: str, max_chars: int) -> tuple[str, bool]:
    """Retire les e-mails (PII), normalise les espaces, plafonne. Renvoie (texte, tronqué)."""
    cleaned = _collapse_ws(_EMAIL_RE.sub("", text))
    if len(cleaned) > max_chars:
        return cleaned[:max_chars].strip(), True
    return cleaned.strip(), False


def _extract_pdf(data: bytes) -> tuple[str, int, int]:
    """Renvoie (texte brut, nb pages, nb images ignorées)."""
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(io.BytesIO(data))
    except (PdfReadError, OSError, ValueError) as exc:
        raise DocumentExtractionError(f"PDF illisible : {exc}") from exc
    if reader.is_encrypted:
        # Tentative de déchiffrement à vide (PDF « protégé » sans mot de passe réel).
        try:
            reader.decrypt("")
        except Exception as exc:  # noqa: BLE001 — pypdf lève des types variés ici
            raise DocumentExtractionError("PDF chiffré (mot de passe requis)") from exc
    parts: list[str] = []
    images = 0
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:  # noqa: BLE001 — page défectueuse : on saute sans casser tout
            continue
        try:
            images += len(page.images)
        except Exception:  # noqa: BLE001 — comptage best-effort
            pass
    return "\n\n".join(parts), len(reader.pages), images


def _extract_docx(data: bytes) -> tuple[str, int]:
    """Renvoie (texte brut paragraphes + tableaux, nb images ignorées)."""
    import docx  # python-docx
    from docx.opc.exceptions import PackageNotFoundError

    try:
        document = docx.Document(io.BytesIO(data))
    except (PackageNotFoundError, KeyError, ValueError, OSError, zipfile.BadZipFile) as exc:
        raise DocumentExtractionError(f"DOCX illisible : {exc}") from exc
    parts = [p.text for p in document.paragraphs if p.text and p.text.strip()]
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    # Images incorporées : parts relationnelles de type image.
    images = sum(
        1 for rel in document.part.rels.values() if "image" in rel.reltype
    )
    return "\n".join(parts), images


def _extract_pptx(data: bytes) -> tuple[str, int, int]:
    """Renvoie (texte des shapes + notes, nb diapositives, nb images ignorées)."""
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError

    try:
        prs = Presentation(io.BytesIO(data))
    except (PackageNotFoundError, KeyError, ValueError, OSError, zipfile.BadZipFile) as exc:
        raise DocumentExtractionError(f"PPTX illisible : {exc}") from exc
    parts: list[str] = []
    images = 0
    slides = 0
    for slide in prs.slides:
        slides += 1
        for shape in slide.shapes:
            # 13 == MSO_SHAPE_TYPE.PICTURE ; on compte sans importer l'enum.
            if getattr(shape, "shape_type", None) == 13:
                images += 1
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
        notes = slide.notes_slide if slide.has_notes_slide else None
        if notes and notes.notes_text_frame and notes.notes_text_frame.text.strip():
            parts.append(notes.notes_text_frame.text)
    return "\n".join(parts), slides, images


def _extract_txt(data: bytes) -> str:
    """Décodage tolérant utf-8 → latin-1 (jamais d'échec sur du texte)."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        return data.decode("latin-1", errors="replace")


def extract_document_text(
    data: bytes, filename: str, *, max_chars: int = DEFAULT_MAX_CHARS
) -> ExtractedDocument:
    """Extrait le texte assaini d'un document joint.

    Args:
        data: contenu binaire du fichier.
        filename: nom d'origine (sert uniquement à déterminer l'extension).
        max_chars: plafond du texte conservé (défaut ``DEFAULT_MAX_CHARS``).

    Raises:
        DocumentExtractionError: format non géré, fichier corrompu ou vide.
    """
    if not data:
        raise DocumentExtractionError("Fichier vide.")
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    if ext not in SUPPORTED_EXTENSIONS:
        raise DocumentExtractionError(
            f"Format non géré : « {ext or filename} ». Formats acceptés : "
            f"{', '.join(SUPPORTED_EXTENSIONS)} (convertissez les .doc/.ppt hérités)."
        )

    warnings: list[str] = []
    if ext == ".pdf":
        raw, pages, images = _extract_pdf(data)
        text, truncated = _sanitize(raw, max_chars)
        if not text:
            warnings.append(
                "Aucun texte extrait — PDF probablement scanné (image). "
                "L'analyse visuelle viendra plus tard."
            )
        return ExtractedDocument(
            text=text, format="pdf", pages=pages, images_skipped=images,
            truncated=truncated, warnings=warnings,
        )
    if ext == ".docx":
        raw, images = _extract_docx(data)
        text, truncated = _sanitize(raw, max_chars)
        if not text:
            warnings.append("Aucun texte extrait du document Word.")
        return ExtractedDocument(
            text=text, format="docx", images_skipped=images,
            truncated=truncated, warnings=warnings,
        )
    if ext == ".pptx":
        raw, slides, images = _extract_pptx(data)
        text, truncated = _sanitize(raw, max_chars)
        if not text:
            warnings.append("Aucun texte extrait des diapositives.")
        return ExtractedDocument(
            text=text, format="pptx", slides=slides, images_skipped=images,
            truncated=truncated, warnings=warnings,
        )
    # .txt
    text, truncated = _sanitize(_extract_txt(data), max_chars)
    if not text:
        raise DocumentExtractionError("Fichier texte vide après nettoyage.")
    return ExtractedDocument(text=text, format="txt", truncated=truncated, warnings=warnings)
