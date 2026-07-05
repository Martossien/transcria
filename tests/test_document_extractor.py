"""Tests de l'extracteur de texte des documents joints (GPU-free, déterministe).

Les fixtures sont fabriquées en mémoire avec les mêmes libs que la production (pypdf,
python-docx, python-pptx) : pas de fichier binaire versionné, pas de donnée réelle.
"""
from __future__ import annotations

import io

import pytest

from transcria.context.document_extractor import (
    DocumentExtractionError,
    ExtractedDocument,
    extract_document_text,
)


def _make_pdf_pypdf(text: str) -> bytes:
    """Fabrique un PDF mono-page contenant ``text`` via pypdf pur (pas de reportlab)."""
    from pypdf import PdfWriter
    from pypdf.generic import (
        ArrayObject,
        DecodedStreamObject,
        DictionaryObject,
        FloatObject,
        NameObject,
        NumberObject,
    )

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    page = writer.pages[0]

    # Flux de contenu minimal traçant du texte avec la police standard Helvetica.
    stream = DecodedStreamObject()
    escaped = text.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
    content = f"BT /F1 12 Tf 10 100 Td ({escaped}) Tj ET".encode("latin-1")
    stream.set_data(content)
    content_ref = writer._add_object(stream)

    font = DictionaryObject()
    font[NameObject("/Type")] = NameObject("/Font")
    font[NameObject("/Subtype")] = NameObject("/Type1")
    font[NameObject("/BaseFont")] = NameObject("/Helvetica")
    font_ref = writer._add_object(font)

    resources = DictionaryObject()
    fonts = DictionaryObject()
    fonts[NameObject("/F1")] = font_ref
    resources[NameObject("/Font")] = fonts
    page[NameObject("/Resources")] = resources
    page[NameObject("/Contents")] = content_ref
    page[NameObject("/MediaBox")] = ArrayObject(
        [NumberObject(0), NumberObject(0), FloatObject(200), FloatObject(200)]
    )

    buf = io.BytesIO()
    writer.write(buf)
    return buf.getvalue()


def _make_docx(paragraphs: list[str], table_rows: list[list[str]] | None = None) -> bytes:
    import docx

    document = docx.Document()
    for p in paragraphs:
        document.add_paragraph(p)
    if table_rows:
        table = document.add_table(rows=len(table_rows), cols=len(table_rows[0]))
        for r, row in enumerate(table_rows):
            for c, val in enumerate(row):
                table.cell(r, c).text = val
    buf = io.BytesIO()
    document.save(buf)
    return buf.getvalue()


def _make_pptx(slides_text: list[list[str]]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    blank = prs.slide_layouts[6]
    for texts in slides_text:
        slide = prs.slides.add_slide(blank)
        for i, t in enumerate(texts):
            box = slide.shapes.add_textbox(Inches(1), Inches(1 + i), Inches(5), Inches(1))
            box.text_frame.text = t
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- PDF ---------------------------------------------------------------------

def test_pdf_extracts_text():
    data = _make_pdf_pypdf("Ordre du jour trimestriel")
    result = extract_document_text(data, "support.pdf")
    assert isinstance(result, ExtractedDocument)
    assert result.format == "pdf"
    assert result.pages == 1
    assert "Ordre du jour" in result.text


def test_pdf_scanned_no_text_warns_not_raises():
    from pypdf import PdfWriter

    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    buf = io.BytesIO()
    writer.write(buf)
    result = extract_document_text(buf.getvalue(), "scan.pdf")
    assert result.text == ""
    assert result.warnings  # message « PDF scanné »


def test_pdf_corrupt_raises():
    with pytest.raises(DocumentExtractionError):
        extract_document_text(b"%PDF-1.4 broken garbage", "bad.pdf")


# --- DOCX --------------------------------------------------------------------

def test_docx_extracts_paragraphs_and_tables():
    data = _make_docx(
        ["Compte rendu de cadrage", "Points à valider"],
        table_rows=[["Action", "Responsable"], ["Livrer le rapport", "Équipe A"]],
    )
    result = extract_document_text(data, "cadrage.docx")
    assert result.format == "docx"
    assert "Compte rendu de cadrage" in result.text
    assert "Livrer le rapport" in result.text
    assert "Responsable" in result.text


def test_docx_corrupt_raises():
    with pytest.raises(DocumentExtractionError):
        extract_document_text(b"not a real docx", "bad.docx")


# --- PPTX --------------------------------------------------------------------

def test_pptx_extracts_slides_and_counts():
    data = _make_pptx([["Titre de la présentation"], ["Point 2", "Détail 2"]])
    result = extract_document_text(data, "deck.pptx")
    assert result.format == "pptx"
    assert result.slides == 2
    assert "Titre de la présentation" in result.text
    assert "Détail 2" in result.text


def test_pptx_corrupt_raises():
    with pytest.raises(DocumentExtractionError):
        extract_document_text(b"PK broken", "bad.pptx")


# --- TXT ---------------------------------------------------------------------

def test_txt_extracts_and_decodes_latin1():
    data = "Réunion très importante".encode("latin-1")
    result = extract_document_text(data, "notes.txt")
    assert result.format == "txt"
    assert "importante" in result.text


def test_txt_empty_raises():
    with pytest.raises(DocumentExtractionError):
        extract_document_text(b"   \n  ", "empty.txt")


# --- Transverse --------------------------------------------------------------

def test_emails_are_stripped_pii():
    data = _make_docx(["Contact : jean.dupont@example.com pour le suivi"])
    result = extract_document_text(data, "c.docx")
    assert "@" not in result.text
    assert "example.com" not in result.text
    assert "pour le suivi" in result.text


def test_truncation_flag_and_cap():
    long_text = "mot " * 10000
    data = _make_docx([long_text])
    result = extract_document_text(data, "long.docx", max_chars=500)
    assert result.truncated is True
    assert len(result.text) <= 500


def test_unsupported_extension_raises():
    with pytest.raises(DocumentExtractionError) as exc:
        extract_document_text(b"whatever", "vieux.ppt")
    assert ".ppt" in str(exc.value) or "non géré" in str(exc.value)


def test_no_extension_raises():
    with pytest.raises(DocumentExtractionError):
        extract_document_text(b"data", "sansextension")


def test_empty_data_raises():
    with pytest.raises(DocumentExtractionError):
        extract_document_text(b"", "x.pdf")
